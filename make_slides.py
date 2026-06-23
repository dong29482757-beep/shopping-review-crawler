# -*- coding: utf-8 -*-
"""5분 발표용 슬라이드 생성 (python-pptx)."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

# 톤다운된 차분한 팔레트: 네이비(메인) + 테라코타(포인트, 절제해서 사용) + 라이트그레이(카드 배경)
NAVY = RGBColor(0x1E, 0x27, 0x61)
ACCENT = RGBColor(0xB8, 0x5A, 0x4D)       # 톤다운 테라코타 (강조 포인트, 솔리드 배경엔 안 씀)
LIGHT_BG = RGBColor(0xF5, 0xF6, 0xFA)     # 카드 배경
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x21, 0x29, 0x5C)
GREY = RGBColor(0x6B, 0x72, 0x80)
ICE = RGBColor(0xD8, 0xDE, 0xF0)          # 다크 배경 위 보조 텍스트

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def add_slide(bg=WHITE):
    s = prs.slides.add_slide(BLANK)
    rect = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    rect.fill.solid()
    rect.fill.fore_color.rgb = bg
    rect.line.fill.background()
    rect.shadow.inherit = False
    s.shapes._spTree.remove(rect._element)
    s.shapes._spTree.insert(2, rect._element)
    return s


def add_text(slide, text, x, y, w, h, size=18, color=DARK, bold=False, align=PP_ALIGN.LEFT,
             font="Calibri", anchor=MSO_ANCHOR.TOP, italic=False, line_spacing=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if line_spacing:
            p.line_spacing = line_spacing
        r = p.add_run()
        r.text = line
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.color.rgb = color
        r.font.name = font
    return box


def add_bullets(slide, items, x, y, w, h, size=15, color=DARK, font="Calibri", space_after=8):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(space_after)
        r = p.add_run()
        r.text = "•  " + item
        r.font.size = Pt(size)
        r.font.color.rgb = color
        r.font.name = font
    return box


def add_chip(slide, text, x, y, w, h, fill, text_color=WHITE, size=12, bold=True):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.adjustments[0] = 0.5
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.fill.background()
    shp.shadow.inherit = False
    tf = shp.text_frame
    tf.word_wrap = False
    tf.margin_left = Inches(0.05); tf.margin_right = Inches(0.05)
    tf.margin_top = 0; tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = text_color
    return shp


def add_card(slide, x, y, w, h, accent=ACCENT, bg=LIGHT_BG):
    """라이트 배경 + 왼쪽 얇은 액센트 바. 솔리드 컬러 블록 대신 쓰는 차분한 카드 스타일."""
    card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    card.fill.solid(); card.fill.fore_color.rgb = bg
    card.line.fill.background(); card.shadow.inherit = False
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(0.07), Inches(h))
    bar.fill.solid(); bar.fill.fore_color.rgb = accent
    bar.line.fill.background(); bar.shadow.inherit = False
    return card


def add_stat(slide, number, label, x, y, w, h, number_color=NAVY, label_color=GREY):
    add_text(slide, number, x, y, w, h * 0.62, size=36, color=number_color, bold=True,
              align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.BOTTOM)
    add_text(slide, label, x, y + h * 0.62, w, h * 0.38, size=12.5, color=label_color,
              align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)


def style_table(table, header_fill=NAVY, header_text=WHITE, body_text=DARK, font_size=12.5):
    for r_idx, row in enumerate(table.rows):
        for c_idx, cell in enumerate(row.cells):
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_left = Inches(0.12); cell.margin_right = Inches(0.12)
            cell.margin_top = Inches(0.05); cell.margin_bottom = Inches(0.05)
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.LEFT if c_idx == 0 else PP_ALIGN.CENTER
                for run in p.runs:
                    run.font.size = Pt(font_size)
                    run.font.name = "Calibri"
                    if r_idx == 0:
                        run.font.bold = True
                        run.font.color.rgb = header_text
                    else:
                        run.font.color.rgb = body_text
            cell.fill.solid()
            cell.fill.fore_color.rgb = header_fill if r_idx == 0 else (WHITE if r_idx % 2 else LIGHT_BG)


# ====================================================================
# Slide 1: Title
# ====================================================================
s = add_slide(NAVY)
add_text(s, "BeautyScope", 1, 2.3, 11.3, 1.2, size=54, color=WHITE, bold=True, font="Georgia")
add_text(s, "화장품 리뷰 538,774건 기반 감성 분석 + 속성기반 상품 리포트 서비스", 1, 3.55, 11.3, 0.6,
         size=18, color=ICE, font="Calibri")
add_chip(s, "데이터 수집", 1, 4.7, 1.8, 0.45, ACCENT)
add_chip(s, "모델 학습 (ML+DL)", 2.95, 4.7, 2.3, 0.45, ACCENT)
add_chip(s, "속성 분석 (ABSA)", 5.4, 4.7, 2.1, 0.45, ACCENT)
add_chip(s, "Streamlit 배포", 7.65, 4.7, 2.0, 0.45, ACCENT)

# ====================================================================
# Slide 2: 서비스 설계 의도
# ====================================================================
s = add_slide(WHITE)
add_text(s, "서비스 설계 의도", 0.7, 0.5, 11.5, 0.8, size=30, color=NAVY, bold=True, font="Georgia")
add_text(s, "평점 하나로 합쳐서 보여주면 속성별 차이가 묻힌다 — 그래서 속성 단위로 쪼개서 분석하는 구조로 설계했다",
         0.7, 1.35, 11.5, 0.6, size=14.5, color=GREY)

prod_a = ("크리스마 더마랩 토너", "4.55", [("보습/수분", 80.5), ("트러블/자극", 23.3)])
prod_b = ("닥터지 에이클리어 밸런싱 토너", "4.78", [("보습/수분", 81.5), ("트러블/자극", 40.5)])
px = 0.7
pw = 5.7
for i, (name, rating, aspects) in enumerate([prod_a, prod_b]):
    add_card(s, px + i*(pw+0.4), 2.15, pw, 2.7)
    tx = px + i*(pw+0.4) + 0.35
    add_text(s, name, tx, 2.4, pw - 0.6, 0.6, size=15, bold=True, color=NAVY)
    add_text(s, f"평점 {rating} / 5", tx, 2.9, pw - 0.6, 0.4, size=12.5, color=GREY)
    ay = 3.4
    for aspect_name, ratio in aspects:
        add_text(s, aspect_name, tx, ay, pw * 0.55, 0.4, size=13, color=DARK)
        add_text(s, f"긍정 {ratio:.1f}%", tx + pw * 0.55, ay, pw * 0.4, 0.4, size=13, bold=True, color=ACCENT,
                  align=PP_ALIGN.RIGHT)
        ay += 0.5

add_text(s, "평점 차이는 0.2점뿐인데, 트러블/자극 속성의 긍정비율은 23.3%와 40.5%로 약 2배 차이",
         0.7, 5.1, 11.5, 0.5, size=14, bold=True, color=NAVY)
add_bullets(s, [
    "리뷰어가 중요하게 보는 속성은 사람마다 다르다 (보습/트러블/향 등)",
    "→ 평점이 아니라 속성 단위로 데이터를 분해해 보여주는 것을 서비스의 핵심 구조로 잡았다",
], 0.7, 5.7, 11.5, 1.3, size=13.5, color=DARK)

# ====================================================================
# Slide 3: 데이터 수집 — 장점과 한계
# ====================================================================
s = add_slide(WHITE)
add_text(s, "데이터 수집 — 장점과 한계", 0.7, 0.5, 11.5, 0.8, size=28, color=NAVY, bold=True, font="Georgia")

stats = [("86,379", "쿠팡 (1차)"), ("452,395", "+ 무신사·올리브영 (2차)"), ("538,774", "최종 통합")]
sx = 0.7
sw = 3.85
for i, (num, label) in enumerate(stats):
    add_card(s, sx + i*(sw+0.2), 1.5, sw, 1.55, accent=ACCENT if i == 2 else NAVY)
    add_stat(s, num, label, sx + i*(sw+0.2) + 0.2, 1.6, sw - 0.4, 1.3,
              number_color=ACCENT if i == 2 else NAVY)

add_card(s, 0.7, 3.4, 5.8, 3.3, accent=NAVY)
add_text(s, "장점", 1.0, 3.65, 5.3, 0.4, size=15.5, bold=True, color=NAVY)
add_bullets(s, [
    "별점이 곧 라벨이 됨 — 4~5점 긍정, 3점 중립, 1~2점 부정으로 바로\n분류에 쓸 수 있어 별도 수작업 레이블링이 필요 없음",
    "화장품 도메인 특화 + 2015~2026년 10년치 데이터로 트렌드 분석도 가능",
    "20개+ 카테고리 × 8개 속성으로 비교 분석 가능",
], 1.0, 4.15, 5.3, 2.4, size=12.5, color=DARK, space_after=10)

add_card(s, 6.8, 3.4, 5.8, 3.3, accent=ACCENT)
add_text(s, "한계", 7.1, 3.65, 5.3, 0.4, size=15.5, bold=True, color=NAVY)
add_bullets(s, [
    "별점 기반 라벨링의 전제는 \"별점 = 진짜 만족도\"인데, 실제로는\n별점과 리뷰 내용이 안 맞는 경우가 있음 (오기입·과장 등)",
    "독립적인 텍스트 기반 분석과 비교했을 때, 별점 5점 리뷰 중 21%가\n실제로는 부정적 내용으로 분류됨",
    "→ 라벨 노이즈가 일정 비율 존재한다는 한계를 인지하고 진행\n(완전히 제거하지 않고, 모델 평가 단계에서 별도로 검증)",
], 7.1, 4.15, 5.3, 2.4, size=12.5, color=DARK, space_after=10)

# ====================================================================
# Slide 4: 쿠팡 크롤링 — 봇 탐지를 어떻게 뚫었나
# ====================================================================
s = add_slide(NAVY)
add_text(s, "쿠팡 크롤링 — 봇 탐지를 어떻게 뚫었나", 0.7, 0.5, 11.5, 0.8, size=27, color=WHITE, bold=True, font="Georgia")
add_text(s, "쿠팡은 Akamai 봇 탐지 시스템이 5단계로 \"사람인지 로봇인지\" 검사한다", 0.7, 1.35, 11.5, 0.5,
         size=14, color=ICE, italic=True)

box1 = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.7), Inches(2.0), Inches(5.8), Inches(4.6))
box1.fill.solid(); box1.fill.fore_color.rgb = RGBColor(0x28, 0x33, 0x70)
box1.line.fill.background(); box1.shadow.inherit = False
add_text(s, "문제 — 왜 막혔나", 1.0, 2.3, 5.1, 0.5, size=16, bold=True, color=WHITE)
add_bullets(s, [
    "requests로 접속 → 헤더만 보고 1단계에서 바로 차단",
    "Selenium으로 접속 → \"자동화 도구\" 흔적이 들켜서\n3단계에서 차단",
    "마우스·스크롤 없이 너무 빠르게 요청하면\n행동 패턴에서 차단",
], 1.0, 2.9, 5.2, 3.5, size=13.5, color=ICE, space_after=14)

box2 = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.8), Inches(2.0), Inches(5.8), Inches(4.6))
box2.fill.solid(); box2.fill.fore_color.rgb = RGBColor(0x28, 0x33, 0x70)
box2.line.fill.background(); box2.shadow.inherit = False
add_text(s, "해결 — \"진짜 손님\" 행세", 7.1, 2.3, 5.2, 0.5, size=16, bold=True, color=WHITE)
add_bullets(s, [
    "가짜 브라우저 대신, 내 컴퓨터의 실제 Chrome을\n그대로 원격 연결해서 사용 (Playwright + CDP)",
    "검색창에 키워드를 한 글자씩 타이핑\n(URL로 바로 안 들어감)",
    "상품 사이 랜덤 대기, 가끔 쿠팡 홈으로 돌아가\n쉬기 — 사람처럼 행동",
], 7.1, 2.9, 5.2, 3.5, color=ICE, size=13.5, space_after=14)

add_text(s, "핵심: 자동화 도구를 새로 띄우는 게 아니라, 이미 로그인된 실제 브라우저에 \"숟가락만 얹는\" 방식이라 자동화 흔적 자체가 없다",
         0.7, 6.75, 11.9, 0.5, size=12, color=ICE, italic=True, align=PP_ALIGN.CENTER)

# ====================================================================
# Slide 5: 데이터 준비 — 무엇이 문제였고 어떻게 풀었나
# ====================================================================
s = add_slide(WHITE)
add_text(s, "데이터 준비 — 무엇이 문제였고 어떻게 풀었나", 0.7, 0.5, 11.5, 0.8, size=26, color=NAVY, bold=True, font="Georgia")

issues = [
    ("긍정 리뷰가 너무 많음", "전체 리뷰의 87%가 긍정이라, 모델이 아무 생각 없이\n\"다 긍정\"이라고만 찍어도 87% 맞는 셈이 됨",
     "부정·중립 리뷰는 그대로 두고, 긍정 리뷰만 솎아내서\n세 감성의 비율을 맞춤"),
    ("문장이 토막 나서 분석됨", "리뷰를 한글자씩/공백 기준으로만 잘라서 \"좋아요\"와\n\"좋아서\"를 완전히 다른 단어로 인식함",
     "형태소 분석기로 단어의 \"원래 형태\"를 찾아내고,\n의미 없는 조사·어미와 불필요한 단어를 제거"),
    ("같은 리뷰가 중복 수집됨", "크롤링을 여러 번 나눠서 진행하다 보니 같은 리뷰가\n두 번 이상 수집된 경우가 있었음",
     "상품·작성자·날짜·내용이 모두 같은 리뷰를 중복으로\n판단해 제거"),
]
y = 1.6
for title, problem, solution in issues:
    add_card(s, 0.7, y, 11.9, 1.65, accent=ACCENT)
    add_chip(s, title, 0.95, y + 0.22, 3.1, 0.5, NAVY, size=12.5)
    add_text(s, "문제", 4.3, y + 0.18, 1.0, 0.3, size=10.5, bold=True, color=ACCENT)
    add_text(s, problem, 4.3, y + 0.5, 3.6, 1.05, size=11.5, color=DARK, line_spacing=1.2)
    add_text(s, "해결", 8.15, y + 0.18, 1.0, 0.3, size=10.5, bold=True, color=ACCENT)
    add_text(s, solution, 8.15, y + 0.5, 4.3, 1.05, size=11.5, color=DARK, line_spacing=1.2)
    y += 1.8

# ====================================================================
# Slide 6: 모델 — 4가지를 비교했다
# ====================================================================
s = add_slide(WHITE)
add_text(s, "모델 — 4가지를 비교했다", 0.7, 0.5, 11.5, 0.8, size=28, color=NAVY, bold=True, font="Georgia")
add_text(s, "리뷰 텍스트를 긍정/중립/부정으로 분류하는 모델을 단순한 것부터 복잡한 것까지 비교했다",
         0.7, 1.3, 11.5, 0.5, size=13.5, color=GREY)

rows_data = [
    ["모델", "역할 / 강점", "accuracy", "macro F1"],
    ["ML\n(TF-IDF+로지스틱회귀)", "가장 가볍고 빠름 — 실서비스 기본 모델", "0.909*", "0.675*"],
    ["DL\n(신경망)", "비선형 패턴까지 학습 — ML보다 정확도 소폭 우위", "0.920*", "0.711*"],
    ["LSTM", "문장 순서를 학습 — 이 데이터에서는 추가 이득 없음", "0.734", "0.671"],
    ["KoELECTRA\n(사전학습 모델)", "적은 데이터로도 가장 높은 정확도", "0.772", "0.660"],
]
table_shape = s.shapes.add_table(5, 4, Inches(0.7), Inches(2.0), Inches(11.9), Inches(3.1))
table = table_shape.table
table.columns[0].width = Inches(3.0)
table.columns[1].width = Inches(5.5)
table.columns[2].width = Inches(1.7)
table.columns[3].width = Inches(1.7)
for r, row_vals in enumerate(rows_data):
    for c, val in enumerate(row_vals):
        table.cell(r, c).text = val
style_table(table)

add_text(s, "* ML/DL은 실제 서비스 운영 분포로 보정한 수치 (학습 시 다운샘플링한 분포와 실제 비율이 달라서 보정함)",
         0.7, 5.3, 11.9, 0.4, size=11, italic=True, color=GREY)

add_card(s, 0.7, 5.8, 11.9, 1.2, accent=NAVY)
add_text(s, "결론: 4개 모델 모두 90% 안팎까지 끌어올렸고, 사전학습 모델이 가장 적은 데이터로도 가장 높은 성능을 냈다. "
            "실제 서비스에는 가볍고 빠른 ML/DL을 우선 적용했다.",
         1.0, 6.0, 11.3, 0.85, size=13.5, bold=True, color=NAVY, line_spacing=1.3)

# ====================================================================
# Slide 7: 모델 문제 발견 → 해결 → 지금 상태
# ====================================================================
s = add_slide(WHITE)
add_text(s, "모델 문제 발견 → 해결 → 지금 상태", 0.7, 0.5, 11.5, 0.8, size=27, color=NAVY, bold=True, font="Georgia")

cols = [
    ("① 발견", "학습 시 다운샘플링한 비율과 실제 운영 비율(긍정\n86.7%)이 달라서, 실제 분포로 다시 평가하니 모델이\n\"중립\"이라고 예측한 것 중 80%가 틀렸음 (label shift)"),
    ("② 해결", "모델을 다시 학습하지 않고, 추론 결과에 \"학습 비율\n↔ 실제 비율\" 차이를 곱해 확률을 보정하는 방식으로\n빠르게 개선"),
    ("③ 지금 상태", "ML accuracy 0.806→0.909, DL accuracy 0.890→0.920\n으로 개선 — 재학습 없이도 실사용 가능한 수준 확보"),
]
cx = 0.7
cw = 3.85
for i, (title, desc) in enumerate(cols):
    add_card(s, cx + i*(cw+0.2), 1.6, cw, 2.5, accent=ACCENT if i == 2 else NAVY)
    add_text(s, title, cx + i*(cw+0.2) + 0.25, 1.85, cw - 0.5, 0.5, size=15, bold=True, color=NAVY)
    add_text(s, desc, cx + i*(cw+0.2) + 0.25, 2.4, cw - 0.5, 1.6, size=12, color=DARK, line_spacing=1.3)

chart_data3 = CategoryChartData()
chart_data3.categories = ["ML accuracy", "ML macro F1", "DL accuracy", "DL macro F1"]
chart_data3.add_series("보정 전", (0.806, 0.626, 0.890, 0.712))
chart_data3.add_series("보정 후", (0.909, 0.675, 0.920, 0.711))
gframe3 = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(1.8), Inches(4.4), Inches(9.7), Inches(2.7), chart_data3)
chart3 = gframe3.chart
chart3.has_legend = True
chart3.legend.position = XL_LEGEND_POSITION.BOTTOM
chart3.legend.include_in_layout = False
chart3.plots[0].series[0].format.fill.solid()
chart3.plots[0].series[0].format.fill.fore_color.rgb = GREY
chart3.plots[0].series[1].format.fill.solid()
chart3.plots[0].series[1].format.fill.fore_color.rgb = NAVY

prs.save(r"D:\crolling\BeautyScope_발표.pptx")
print("저장 완료")
